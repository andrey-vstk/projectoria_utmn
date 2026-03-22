from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BRAND = RGBColor(0, 174, 239)
BRAND_DARK = RGBColor(0, 122, 170)
TEXT = RGBColor(26, 32, 44)
MUTED = RGBColor(91, 104, 121)
LIGHT_BG = RGBColor(242, 251, 255)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Layout helpers

def add_base(slide, show_top_bar=True, footer='Проектория • ТюмГУ • 2026'):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    if show_top_bar:
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            prs.slide_width,
            Inches(0.28),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = BRAND
        top_bar.line.fill.background()

    footer_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(7.08), Inches(12.3), Inches(0.3)
    )
    tf = footer_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = footer
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(0.65), Inches(11.9), Inches(1.5))
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Helvetica Neue'
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = TEXT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = 'Helvetica Neue'
        p2.font.size = Pt(17)
        p2.font.color.rgb = MUTED


def add_bullets(slide, items, x=0.95, y=1.9, w=7.8, h=4.8, font_size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = 'Helvetica Neue'
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)


def add_highlight_card(slide, title, lines, x=8.9, y=1.8, w=3.8, h=4.8):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = BRAND

    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.22), Inches(w - 0.44), Inches(0.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = 'Helvetica Neue'
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = BRAND_DARK

    content = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.85), Inches(w - 0.44), Inches(h - 1.0))
    ctf = content.text_frame
    ctf.clear()
    for i, line in enumerate(lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = f'• {line}'
        p.font.name = 'Helvetica Neue'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def add_metric_card(slide, title, value, x, y, w=3.0, h=1.5):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = BRAND

    tb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.2), Inches(w - 0.44), Inches(h - 0.4))
    tf = tb.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = 'Helvetica Neue'
    p1.font.size = Pt(13)
    p1.font.color.rgb = MUTED

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Helvetica Neue'
    p2.font.bold = True
    p2.font.size = Pt(27)
    p2.font.color.rgb = BRAND_DARK


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide, show_top_bar=False)

left_strip = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1.35), prs.slide_height
)
left_strip.fill.solid()
left_strip.fill.fore_color.rgb = BRAND
left_strip.line.fill.background()

add_title(
    slide,
    'Проектория',
    'ВКР: «Проектирование веб-системы поддержки формирования межподразделенческих\nуниверситетских проектов по текстовым запросам индустрии»',
)

subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.2), Inches(2.0))
stf = subtitle_box.text_frame
stf.clear()
p = stf.paragraphs[0]
p.text = 'Результат практики: полнофункциональный MVP разработан с нуля'
p.font.name = 'Helvetica Neue'
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = BRAND_DARK

p2 = stf.add_paragraph()
p2.text = 'Студент: ____________________      Руководитель: ____________________'
p2.font.name = 'Helvetica Neue'
p2.font.size = Pt(17)
p2.font.color.rgb = MUTED
p2.space_before = Pt(18)

# Slide 2: Problem and relevance
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Актуальность и проблема')
add_bullets(
    slide,
    [
        'Запросы индустрии приходят в виде неструктурированных стенограмм и описаний.',
        'Маршрутизация между подразделениями часто выполняется вручную.',
        'Теряется контекст, увеличивается время реакции университета.',
        'Нет прозрачной воронки: от запроса до отклика подразделений.',
    ],
)
add_highlight_card(
    slide,
    'Ценность решения',
    [
        'Ускорение отбора релевантных направлений',
        'Единая цифровая история по каждому проекту',
        'Повышение управляемости межподразделенческого взаимодействия',
    ],
)

# Slide 3: Goal and tasks
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Цель и задачи практики')
add_bullets(
    slide,
    [
        'Цель: спроектировать и реализовать MVP веб-системы для запуска межподразделенческих проектов.',
        'Задача 1: построить безопасный веб-сервис с ролевой моделью.',
        'Задача 2: реализовать асинхронный LLM-пайплайн с контролем статусов.',
        'Задача 3: организовать подтверждаемую рассылку и фиксацию откликов.',
        'Задача 4: обеспечить админ-управление пользователями и подразделениями.',
    ],
    w=11.6,
    h=4.9,
    font_size=20,
)

# Slide 4: What was built
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Что реализовано с нуля')

add_highlight_card(
    slide,
    'Backend',
    [
        'NestJS API + Prisma + PostgreSQL',
        'JWT в httpOnly cookie, RBAC (admin/initiator)',
        'BullMQ + Redis для фоновых задач',
        'SMTP-рассылка и публичный endpoint отклика',
        'Интеграция с n8n (webhook + REST контракт)',
    ],
    x=0.95,
    y=1.9,
    w=5.95,
    h=4.8,
)

add_highlight_card(
    slide,
    'Frontend',
    [
        'Next.js интерфейс на русском языке',
        'Страницы: логин, проекты, карточка проекта, отклики, админка',
        'Ручная проверка/редактирование писем перед отправкой',
        'Дизайн-система в фирменном стиле ТюмГУ (#00AEEF)',
    ],
    x=6.95,
    y=1.9,
    w=5.4,
    h=4.8,
)

# Slide 5: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Архитектура MVP')

# Row 1
for label, x in [('Web (Next.js)', 0.95), ('API (NestJS)', 4.35), ('n8n + Ollama', 7.75), ('SMTP', 11.0)]:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.05), Inches(2.15), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BRAND
    tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(2.33), Inches(1.9), Inches(0.5))
    tf = tb.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT

for label, x in [('PostgreSQL', 3.0), ('Redis/BullMQ', 6.0), ('Хранилище\nсущностей', 9.0)]:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.25), Inches(2.25), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BRAND_DARK
    tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(4.53), Inches(2.0), Inches(0.6))
    tf = tb.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Helvetica Neue'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT

flow = slide.shapes.add_textbox(Inches(0.95), Inches(5.75), Inches(11.8), Inches(0.8))
ft = flow.text_frame
ft.text = 'Поток: пользователь → API → очередь → LLM/n8n → результат → ручное подтверждение → рассылка → отклики'
pp = ft.paragraphs[0]
pp.font.name = 'Helvetica Neue'
pp.font.size = Pt(16)
pp.font.color.rgb = BRAND_DARK
pp.font.bold = True

# Slide 6: Scenario
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Сквозной бизнес-сценарий')

steps = [
    '1. Инициатор создает проект и загружает текст',
    '2. Запускается анализ (асинхронная задача)',
    '3. LLM формирует summary, задачи и письма по подразделениям',
    '4. Инициатор проверяет и редактирует предложения',
    '5. Подтверждает рассылку и отправку писем',
    '6. Подразделения откликаются по уникальной ссылке',
    '7. Инициатор получает уведомление и видит отклики в карточке',
]
add_bullets(slide, steps, w=11.8, h=4.9, font_size=20)

# Slide 7: Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Итоги практики')

add_metric_card(slide, 'Собранных сервисов', '8', x=0.95, y=1.9)
add_metric_card(slide, 'Ключевых сущностей БД', '10+', x=4.2, y=1.9)
add_metric_card(slide, 'Ролевых сценариев', '2', x=7.45, y=1.9)
add_metric_card(slide, 'Сквозных потоков', '3', x=10.7, y=1.9)

add_bullets(
    slide,
    [
        'Разработан production-minded каркас: backend, frontend, БД, очереди, интеграции, Docker.',
        'Реализована прозрачная история: запрос → анализ → рассылка → отклики.',
        'Подготовлен фундамент для масштабирования в рамках ВКР.',
    ],
    y=3.9,
    w=11.8,
    h=2.5,
    font_size=20,
)

# Slide 8: Novelty and development
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide)
add_title(slide, 'Новизна и развитие в ВКР')
add_bullets(
    slide,
    [
        'Комбинация LLM-декомпозиции и управляемого межподразделенческого роутинга в единой системе.',
        'Человек в контуре: обязательная ручная валидация перед коммуникацией.',
        'Готовность к локальной модели: замена провайдера без смены бизнес-логики.',
        'Следующий этап ВКР: метрики качества роутинга, оценка точности и производительности.',
        'План развития: SSO/LDAP, корпоративный почтовый шлюз, аналитика воронки, SLA-мониторинг.',
    ],
    w=11.8,
    h=4.9,
    font_size=19,
)

# Slide 9: Final
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_base(slide, show_top_bar=False)

bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = BRAND
bg.line.fill.background()

thank = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.8), Inches(2.5))
tf = thank.text_frame
tf.clear()
p1 = tf.paragraphs[0]
p1.text = 'Спасибо за внимание'
p1.font.name = 'Helvetica Neue'
p1.font.bold = True
p1.font.size = Pt(52)
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = 'Готовы к демонстрации «Проектории» и ответам на вопросы'
p2.font.name = 'Helvetica Neue'
p2.font.size = Pt(22)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(14)

out_path = 'docs/presentation/proektoriya_vkr_defense.pptx'
prs.save(out_path)
print(out_path)
